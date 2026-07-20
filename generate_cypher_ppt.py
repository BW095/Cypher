#!/usr/bin/env python3
"""
generate_cypher_ppt.py — Generate the Cypher project presentation (.pptx)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ─── Colour palette (Formal, White Background) ──────────────────────────
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY    = RGBColor(0xF1, 0xF5, 0xF9)
MID_GRAY      = RGBColor(0x64, 0x74, 0x8B)
DARK_GRAY     = RGBColor(0x33, 0x41, 0x55)
BLACK         = RGBColor(0x0F, 0x17, 0x2A)
ACCENT_BLUE   = RGBColor(0x1E, 0x40, 0xAF)   # Deep Blue
ACCENT_TEAL   = RGBColor(0x0F, 0x76, 0x6E)   # Deep Teal
ACCENT_RED    = RGBColor(0xB9, 0x1C, 0x1C)   # Deep Red for problems
ACCENT_AMBER  = RGBColor(0xB4, 0x53, 0x09)   # Amber

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

BLANK_LAYOUT = prs.slide_layouts[6]

# ─── helper functions ──────────────────────────────────────────────────
def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, border_color=None, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf

def set_para(tf, text, font_size=14, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT, space_after=Pt(6), space_before=Pt(0), font_name="Calibri"):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = space_after
    p.space_before = space_before
    return p

def first_para(tf, text, font_size=14, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT, space_after=Pt(6), font_name="Calibri"):
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = space_after
    return p

def add_divider(slide, left, top, width, color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_section_number(slide, number, left, top, color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(0.7), Inches(0.7))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(4)

def add_slide_number(slide, num, total):
    tf = add_text_box(slide, Inches(11.5), Inches(7.0), Inches(1.5), Inches(0.4))
    first_para(tf, f"{num} / {total}", font_size=10, color=MID_GRAY, alignment=PP_ALIGN.RIGHT)

def add_section_header_slide(slide, section_title, subtitle, section_num, accent_color=ACCENT_BLUE):
    add_bg(slide, WHITE)
    
    # Large section number
    tf = add_text_box(slide, Inches(0.8), Inches(1.2), Inches(3), Inches(1.8))
    first_para(tf, f"0{section_num}", font_size=80, bold=True, color=accent_color, font_name="Calibri")
    
    add_divider(slide, Inches(0.8), Inches(3.3), Inches(2), accent_color)
    
    tf = add_text_box(slide, Inches(0.8), Inches(3.6), Inches(11), Inches(1.5))
    first_para(tf, section_title, font_size=44, bold=True, color=BLACK, font_name="Calibri")
    
    tf = add_text_box(slide, Inches(0.8), Inches(5.2), Inches(10), Inches(1.2))
    first_para(tf, subtitle, font_size=18, color=MID_GRAY, font_name="Calibri")
    
    add_shape(slide, Inches(12.8), Inches(0), Inches(0.533), SLIDE_H, accent_color)

# ═══════════════════════════════════════════════════════════════════════
TOTAL_SLIDES = 25
slide_num = 0

def next_slide():
    global slide_num
    slide_num += 1
    return prs.slides.add_slide(BLANK_LAYOUT)

# ═══════════════════════════════════════════════════════════════════════
# TITLE SLIDE
# ═══════════════════════════════════════════════════════════════════════
slide = next_slide()
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_BLUE)

tf = add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.6))
first_para(tf, "INDUSTRIAL KNOWLEDGE INTELLIGENCE", font_size=14, bold=True, color=ACCENT_BLUE, font_name="Calibri")

tf = add_text_box(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(2.0))
first_para(tf, "CYPHER", font_size=72, bold=True, color=BLACK, font_name="Calibri")
set_para(tf, "Unified Asset & Operations Brain", font_size=36, bold=False, color=ACCENT_TEAL, font_name="Calibri")

add_divider(slide, Inches(0.8), Inches(4.2), Inches(4), ACCENT_BLUE)

tf = add_text_box(slide, Inches(0.8), Inches(4.6), Inches(10), Inches(1.5))
first_para(tf, "An AI-powered platform transforming heterogeneous documents into actionable knowledge.", font_size=22, color=DARK_GRAY, font_name="Calibri")
set_para(tf, "Overcoming the industrial knowledge cliff with Edge-AI, RAG, and Graph computing.", font_size=16, color=MID_GRAY, font_name="Calibri", space_after=Pt(12))

badge = add_shape(slide, Inches(0.8), Inches(6.2), Inches(5.5), Inches(0.5), LIGHT_GRAY, ACCENT_BLUE)
tf_badge = badge.text_frame
tf_badge.word_wrap = True
p = tf_badge.paragraphs[0]
p.text = "Document Management  ·  Knowledge Engineering  ·  Quality"
p.font.size = Pt(11)
p.font.color.rgb = ACCENT_BLUE
p.font.name = "Calibri"
p.alignment = PP_ALIGN.CENTER

add_slide_number(slide, slide_num, TOTAL_SLIDES)

# ═══════════════════════════════════════════════════════════════════════
# SECTION 1 — ABOUT TEAM
# ═══════════════════════════════════════════════════════════════════════
slide = next_slide()
add_section_header_slide(slide, "About The Team", "The creators behind the Cypher knowledge intelligence platform.", 1, ACCENT_BLUE)
add_slide_number(slide, slide_num, TOTAL_SLIDES)

# Slide 1: Team Overview
slide = next_slide()
add_bg(slide, WHITE)
tf = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.8))
first_para(tf, "Team Overview", font_size=36, bold=True, color=BLACK)
add_divider(slide, Inches(0.8), Inches(1.2), Inches(2), ACCENT_BLUE)

member_data = [
    ("Team Member 1", "AI & Backend Lead", "[ Replace with Photo ]\n\nExpertise in Edge-AI deployment, LLMs, and asynchronous API architecture."),
    ("Team Member 2", "Graph & Frontend Lead", "[ Replace with Photo ]\n\nExpertise in Knowledge Graphs (Neo4j), React UI, and data visualization.")
]
for i, (name, role, desc) in enumerate(member_data):
    left = Inches(1.5 + i * 5.5)
    card = add_shape(slide, left, Inches(2.2), Inches(4.5), Inches(4.5), LIGHT_GRAY, ACCENT_BLUE)
    
    photo = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(1.5), Inches(2.6), Inches(1.5), Inches(1.5))
    photo.fill.solid()
    photo.fill.fore_color.rgb = WHITE
    photo.line.color.rgb = ACCENT_BLUE
    photo.line.width = Pt(1.5)
    ph_tf = photo.text_frame
    ph_tf.paragraphs[0].text = "📷"
    ph_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    ph_tf.paragraphs[0].font.size = Pt(28)
    
    tf = add_text_box(slide, left, Inches(4.3), Inches(4.5), Inches(0.5))
    first_para(tf, name, font_size=20, bold=True, color=BLACK, alignment=PP_ALIGN.CENTER)
    
    tf = add_text_box(slide, left, Inches(4.8), Inches(4.5), Inches(0.4))
    first_para(tf, role, font_size=16, bold=True, color=ACCENT_TEAL, alignment=PP_ALIGN.CENTER)
    
    tf = add_text_box(slide, left + Inches(0.5), Inches(5.3), Inches(3.5), Inches(1.0))
    first_para(tf, desc, font_size=12, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, slide_num, TOTAL_SLIDES)

# Slide 2: Division of Expertise
slide = next_slide()
add_bg(slide, WHITE)
tf = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.8))
first_para(tf, "Division of Expertise", font_size=36, bold=True, color=BLACK)
add_divider(slide, Inches(0.8), Inches(1.2), Inches(2), ACCENT_BLUE)

contributions = [
    ("Member 1 Focus", "• Architected FastAPI asynchronous backend\n• Integrated Llama.cpp and Qwen VLM models\n• Developed vector retrieval logic (Qdrant)\n• Implemented document ingestion pipeline", ACCENT_TEAL),
    ("Member 2 Focus", "• Designed Knowledge Graph schema (Neo4j)\n• Built the React dashboard and UI\n• Developed graph-based entity visualization\n• Handled deployment and Docker orchestration", ACCENT_BLUE),
]

for i, (area, desc, color) in enumerate(contributions):
    left = Inches(0.8 + i * 6.1)
    card = add_shape(slide, left, Inches(2.5), Inches(5.8), Inches(3.0), LIGHT_GRAY, color)
    
    tf = add_text_box(slide, left + Inches(0.5), Inches(2.8), Inches(5.0), Inches(0.6))
    first_para(tf, area, font_size=20, bold=True, color=color)
    
    tf = add_text_box(slide, left + Inches(0.5), Inches(3.5), Inches(5.0), Inches(1.8))
    for line in desc.split("\n"):
        set_para(tf, line, font_size=14, color=DARK_GRAY)

add_slide_number(slide, slide_num, TOTAL_SLIDES)

# Slide 3: Vision & Mission
slide = next_slide()
add_bg(slide, WHITE)
tf = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.8))
first_para(tf, "Vision & Mission", font_size=36, bold=True, color=BLACK)
add_divider(slide, Inches(0.8), Inches(1.2), Inches(2), ACCENT_BLUE)

card = add_shape(slide, Inches(0.8), Inches(2.5), Inches(11.5), Inches(2.5), LIGHT_GRAY, ACCENT_TEAL)
tf = add_text_box(slide, Inches(1.2), Inches(2.8), Inches(10.8), Inches(1.6))
first_para(tf, "Our Mission", font_size=24, bold=True, color=ACCENT_TEAL)
set_para(tf, "To solve the industrial knowledge fragmentation crisis by empowering manufacturing and energy companies with an autonomous AI copilot that democratizes access to technical documentation, reducing downtime, and securing institutional knowledge for the future.", font_size=18, color=DARK_GRAY, space_before=Pt(12))

add_slide_number(slide, slide_num, TOTAL_SLIDES)

# ═══════════════════════════════════════════════════════════════════════
# SECTION 2 — PROBLEM STATEMENT
# ═══════════════════════════════════════════════════════════════════════
slide = next_slide()
add_section_header_slide(slide, "Problem Statement", "The cost of fragmented industrial knowledge and the incoming knowledge cliff.", 2, ACCENT_RED)
add_slide_number(slide, slide_num, TOTAL_SLIDES)

# Slide 1: The Industrial Knowledge Crisis
slide = next_slide()
add_bg(slide, WHITE)
tf = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8))
first_para(tf, "The Industrial Knowledge Crisis", font_size=36, bold=True, color=BLACK)
add_divider(slide, Inches(0.8), Inches(1.2), Inches(2), ACCENT_RED)

tf = add_text_box(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.0))
first_para(tf, "Information Fragmentation", font_size=20, bold=True, color=ACCENT_RED)
set_para(tf, "Large plants operate across 7 to 12 disconnected document systems. P&IDs in one place, maintenance orders in another, inspection records in a third.", font_size=16, color=DARK_GRAY, space_before=Pt(8))
set_para(tf, "A 2024 McKinsey global survey found that professionals in asset-intensive industries spend a massive portion of their week just looking for data.", font_size=16, color=DARK_GRAY, space_before=Pt(8))

card = add_shape(slide, Inches(7.0), Inches(2.5), Inches(4.5), Inches(2.5), LIGHT_GRAY, ACCENT_RED)
tf = add_text_box(slide, Inches(7.2), Inches(3.0), Inches(4.1), Inches(1.0))
first_para(tf, "35%", font_size=60, bold=True, color=ACCENT_RED, alignment=PP_ALIGN.CENTER)
tf = add_text_box(slide, Inches(7.2), Inches(4.0), Inches(4.1), Inches(0.5))
first_para(tf, "of working hours spent searching for info or clarifying instructions", font_size=14, bold=True, color=BLACK, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, slide_num, TOTAL_SLIDES)

# Slide 2: The Knowledge Cliff
slide = next_slide()
add_bg(slide, WHITE)
tf = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8))
first_para(tf, "The Knowledge Cliff", font_size=36, bold=True, color=BLACK)
add_divider(slide, Inches(0.8), Inches(1.2), Inches(2), ACCENT_AMBER)

tf = add_text_box(slide, Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.5))
first_para(tf, "An estimated 25% of India's experienced industrial engineers and operators will retire within the next decade.", font_size=20, bold=True, color=BLACK)
set_para(tf, "They will take decades of undocumented operational knowledge, intuition, and failure pattern recognition with them. Once this knowledge leaves the factory floor, it cannot be easily recovered.", font_size=16, color=DARK_GRAY, space_before=Pt(10))

card = add_shape(slide, Inches(0.8), Inches(4.0), Inches(11.5), Inches(2.0), LIGHT_GRAY, ACCENT_AMBER)
tf = add_text_box(slide, Inches(1.2), Inches(4.3), Inches(10.8), Inches(1.5))
first_para(tf, "The Consequence", font_size=18, bold=True, color=ACCENT_AMBER)
set_para(tf, "New hires rely on outdated or scattered manuals. Troubleshooting that took a veteran 10 minutes takes a new operator hours because the context is lost.", font_size=16, color=DARK_GRAY, space_before=Pt(8))

add_slide_number(slide, slide_num, TOTAL_SLIDES)

# Slide 3: Impact of Fragmentation
slide = next_slide()
add_bg(slide, WHITE)
tf = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8))
first_para(tf, "Impact of Fragmentation", font_size=36, bold=True, color=BLACK)
add_divider(slide, Inches(0.8), Inches(1.2), Inches(2), ACCENT_RED)

tf = add_text_box(slide, Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.0))
first_para(tf, "Knowledge fragmentation is not just a file management problem.", font_size=24, bold=True, color=ACCENT_RED)

impacts = [
    ("Safety Risks", "Decisions made without full equipment history or hazard warnings.", ACCENT_AMBER),
    ("Quality Drops", "Inconsistent adherence to evolving regulatory requirements.", ACCENT_TEAL),
    ("Downtime", "18–22% of unplanned downtime events in heavy industry are directly attributed to disconnected data.", ACCENT_RED)
]

for i, (title, desc, color) in enumerate(impacts):
    left = Inches(0.8 + i * 4.0)
    card = add_shape(slide, left, Inches(3.5), Inches(3.6), Inches(2.2), LIGHT_GRAY, color)
    tf = add_text_box(slide, left + Inches(0.2), Inches(3.8), Inches(3.2), Inches(0.5))
    first_para(tf, title, font_size=18, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    tf = add_text_box(slide, left + Inches(0.2), Inches(4.4), Inches(3.2), Inches(1.0))
    first_para(tf, desc, font_size=14, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, slide_num, TOTAL_SLIDES)

# Slide 4: The Core Challenge
slide = next_slide()
add_bg(slide, WHITE)
tf = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8))
first_para(tf, "The Core Challenge", font_size=36, bold=True, color=BLACK)
add_divider(slide, Inches(0.8), Inches(1.2), Inches(2), ACCENT_BLUE)

card = add_shape(slide, Inches(0.8), Inches(2.0), Inches(11.5), Inches(3.0), LIGHT_GRAY, ACCENT_BLUE)
tf = add_text_box(slide, Inches(1.2), Inches(2.5), Inches(10.8), Inches(2.0))
first_para(tf, "Challenge Statement:", font_size=22, bold=True, color=ACCENT_BLUE)
set_para(tf, "Build an AI-powered Industrial Knowledge Intelligence platform that ingests heterogeneous documents — engineering drawings, maintenance records, safety procedures, operating instructions — across structured and unstructured formats.", font_size=18, color=BLACK, space_before=Pt(12))
set_para(tf, "Make their collective intelligence queryable, actionable, and continuously updated at the point of need.", font_size=18, color=DARK_GRAY, space_before=Pt(12))

add_slide_number(slide, slide_num, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════
# SECTION 3 — ARCHITECTURE
# ═══════════════════════════════════════════════════════════════════════
slide = next_slide()
add_section_header_slide(slide, "System Architecture", "The foundations of Cypher: Local LLMs, Vector DB, and Knowledge Graph.", 3, ACCENT_BLUE)
add_slide_number(slide, slide_num, TOTAL_SLIDES)

# Slide 1: System Architecture — Four Pillars
slide = next_slide()
add_bg(slide, WHITE)
tf = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8))
first_para(tf, "System Architecture — Four Pillars", font_size=36, bold=True, color=BLACK)
add_divider(slide, Inches(0.8), Inches(1.2), Inches(2), ACCENT_BLUE)

pillars = [
    ("Neo4j", "Graph Database", "Stores entity relationships, equipment hierarchies, and topological connections.", ACCENT_TEAL),
    ("Qdrant", "Vector Database", "Stores high-dimensional embeddings for semantic search over unstructured text.", ACCENT_BLUE),
    ("FastAPI + LLM", "The Brain Backend", "Orchestrates ingestion, querying, local Llama.cpp models, and BGE embeddings.", ACCENT_AMBER),
    ("React UI", "Frontend Dashboard", "Provides user search, graph visualization, and chat interfaces.", DARK_GRAY)
]

for i, (title, sub, desc, color) in enumerate(pillars):
    left = Inches(0.8 + i * 3.0)
    card = add_shape(slide, left, Inches(2.5), Inches(2.8), Inches(3.5), LIGHT_GRAY, color)
    tf = add_text_box(slide, left + Inches(0.2), Inches(2.8), Inches(2.4), Inches(0.5))
    first_para(tf, title, font_size=20, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    tf = add_text_box(slide, left + Inches(0.2), Inches(3.3), Inches(2.4), Inches(0.3))
    first_para(tf, sub, font_size=14, bold=True, color=BLACK, alignment=PP_ALIGN.CENTER)
    tf = add_text_box(slide, left + Inches(0.2), Inches(3.9), Inches(2.4), Inches(1.5))
    first_para(tf, desc, font_size=12, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, slide_num, TOTAL_SLIDES)

# Slide 2: System Architecture Diagram
slide = next_slide()
add_bg(slide, WHITE)
tf = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8))
first_para(tf, "System Architecture Flow", font_size=36, bold=True, color=BLACK)
add_divider(slide, Inches(0.8), Inches(1.2), Inches(2), ACCENT_BLUE)

img_path = "/home/bw/.gemini/antigravity-ide/brain/24e47e86-2966-44e2-9d30-384023edd237/cypher_architecture_light_1784489745756.png"
# Center the image horizontally
slide.shapes.add_picture(img_path, Inches(3.5), Inches(1.5), width=Inches(6.0), height=Inches(5.7))

add_slide_number(slide, slide_num, TOTAL_SLIDES)

# Slide 3: Data Ingestion & Processing
slide = next_slide()
add_bg(slide, WHITE)
tf = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8))
first_para(tf, "Data Ingestion Pipeline", font_size=36, bold=True, color=BLACK)
add_divider(slide, Inches(0.8), Inches(1.2), Inches(2), ACCENT_TEAL)

steps = [
    ("1. Read & OCR", "Raw PDFs, Images, and Text are loaded via `/documents` mount. Visual models (Qwen3VL) parse P&IDs."),
    ("2. Entity Extraction", "Qwen2.5-3B model extracts key tags, dates, equipment IDs, and specifications from the text chunks."),
    ("3. Vectorization", "BGE embedding model converts text into semantic vectors for Qdrant storage."),
    ("4. Graph Linkage", "Extracted entities are mapped as nodes in Neo4j, establishing relationships between documents and assets.")
]

for i, (title, desc) in enumerate(steps):
    top = Inches(2.0 + i * 1.2)
    card = add_shape(slide, Inches(0.8), top, Inches(11.5), Inches(1.0), LIGHT_GRAY, ACCENT_TEAL)
    tf = add_text_box(slide, Inches(1.2), top + Inches(0.15), Inches(2.5), Inches(0.7))
    first_para(tf, title, font_size=16, bold=True, color=ACCENT_TEAL)
    tf = add_text_box(slide, Inches(3.8), top + Inches(0.15), Inches(8.2), Inches(0.7))
    first_para(tf, desc, font_size=14, color=DARK_GRAY)

add_slide_number(slide, slide_num, TOTAL_SLIDES)

# Slide 3: The Hybrid Decision Engine
slide = next_slide()
add_bg(slide, WHITE)
tf = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8))
first_para(tf, "The Hybrid Retrieval Engine", font_size=36, bold=True, color=BLACK)
add_divider(slide, Inches(0.8), Inches(1.2), Inches(2), ACCENT_BLUE)

tf = add_text_box(slide, Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.0))
first_para(tf, "Semantic Search alone is insufficient for engineering data. We fuse Vector and Graph retrieval.", font_size=18, bold=True, color=BLACK)

card1 = add_shape(slide, Inches(0.8), Inches(3.2), Inches(5.5), Inches(2.5), LIGHT_GRAY, ACCENT_BLUE)
tf = add_text_box(slide, Inches(1.2), Inches(3.5), Inches(4.8), Inches(0.5))
first_para(tf, "Vector Search (Qdrant)", font_size=18, bold=True, color=ACCENT_BLUE)
tf = add_text_box(slide, Inches(1.2), Inches(4.0), Inches(4.8), Inches(1.5))
first_para(tf, "Finds conceptually similar text. If a user asks 'How to fix hydraulic leak?', it retrieves manual sections detailing hydraulic repair procedures.", font_size=14, color=DARK_GRAY)

card2 = add_shape(slide, Inches(6.8), Inches(3.2), Inches(5.5), Inches(2.5), LIGHT_GRAY, ACCENT_TEAL)
tf = add_text_box(slide, Inches(7.2), Inches(3.5), Inches(4.8), Inches(0.5))
first_para(tf, "Graph Search (Neo4j)", font_size=18, bold=True, color=ACCENT_TEAL)
tf = add_text_box(slide, Inches(7.2), Inches(4.0), Inches(4.8), Inches(1.5))
first_para(tf, "Understands structural links. It knows that 'Pump A' is connected to 'Valve B', and retrieves the last 3 maintenance records explicitly linked to 'Pump A'.", font_size=14, color=DARK_GRAY)

add_slide_number(slide, slide_num, TOTAL_SLIDES)

# Slide 4: Local Edge-AI Deployment
slide = next_slide()
add_bg(slide, WHITE)
tf = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8))
first_para(tf, "Local Edge-AI & Privacy", font_size=36, bold=True, color=BLACK)
add_divider(slide, Inches(0.8), Inches(1.2), Inches(2), ACCENT_RED)

tf = add_text_box(slide, Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.0))
first_para(tf, "Industrial data is highly sensitive. Cypher runs 100% locally.", font_size=20, bold=True, color=ACCENT_RED)

features = [
    ("Zero Cloud Dependency", "No proprietary documents, P&IDs, or intellectual property ever leave the facility's internal network."),
    ("Llama.cpp Optimization", "Runs massive quantized GGUF models efficiently on local hardware. The Qwen 8B model fits gracefully into VRAM."),
    ("Hardware Acceleration", "Utilizes the NVIDIA Container Toolkit to pass host GPUs into the Dockerized backend for instant LLM inference."),
]

for i, (title, desc) in enumerate(features):
    top = Inches(3.2 + i * 1.1)
    tf = add_text_box(slide, Inches(1.0), top, Inches(11.0), Inches(0.8))
    first_para(tf, f"• {title}:", font_size=18, bold=True, color=BLACK)
    set_para(tf, desc, font_size=16, color=DARK_GRAY)

add_slide_number(slide, slide_num, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════
# SECTION 4 — TECHNICAL IMPLEMENTATION
# ═══════════════════════════════════════════════════════════════════════
slide = next_slide()
add_section_header_slide(slide, "Technical Implementation", "Deep dive into RAG, AI Pipelines, and asynchronous API design.", 4, ACCENT_AMBER)
add_slide_number(slide, slide_num, TOTAL_SLIDES)

# Slide 1: Document Processing
slide = next_slide()
add_bg(slide, WHITE)
tf = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8))
first_para(tf, "Advanced Document Processing", font_size=36, bold=True, color=BLACK)
add_divider(slide, Inches(0.8), Inches(1.2), Inches(2), ACCENT_AMBER)

tf = add_text_box(slide, Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.5))
first_para(tf, "We use specialized models for varying data modalities:", font_size=18, bold=True, color=BLACK)

card = add_shape(slide, Inches(0.8), Inches(3.0), Inches(11.5), Inches(3.0), LIGHT_GRAY, ACCENT_AMBER)
tf = add_text_box(slide, Inches(1.2), Inches(3.2), Inches(10.8), Inches(0.5))
first_para(tf, "Vision-Language Processing (Qwen3VL-8B-Instruct)", font_size=18, bold=True, color=ACCENT_AMBER)
tf = add_text_box(slide, Inches(1.2), Inches(3.8), Inches(10.8), Inches(2.0))
first_para(tf, "• Translates complex engineering drawings (P&IDs) into textual logic.", font_size=14, color=DARK_GRAY)
set_para(tf, "• Reads tabular data and maintenance spreadsheets.", font_size=14, color=DARK_GRAY)
set_para(tf, "• Bridges the gap between scanned, unstructured imagery and the queryable vector database.", font_size=14, color=DARK_GRAY)
set_para(tf, "• Employs a custom Vision Projector (`mmproj-Qwen3VL`) for highly accurate visual tokenization.", font_size=14, color=DARK_GRAY)

add_slide_number(slide, slide_num, TOTAL_SLIDES)

# Slide 2: RAG Pipeline
slide = next_slide()
add_bg(slide, WHITE)
tf = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8))
first_para(tf, "Expert Knowledge Copilot (RAG)", font_size=36, bold=True, color=BLACK)
add_divider(slide, Inches(0.8), Inches(1.2), Inches(2), ACCENT_BLUE)

tf = add_text_box(slide, Inches(0.8), Inches(2.0), Inches(11.5), Inches(0.5))
first_para(tf, "The Retrieval-Augmented Generation (RAG) Flow:", font_size=18, bold=True, color=BLACK)

flow = [
    ("User Query", "Engineer asks a natural language question via the React UI.", ACCENT_TEAL),
    ("Search Execution", "Backend runs parallel queries against Qdrant (vectors) and Neo4j (graph).", ACCENT_BLUE),
    ("Reranking", "A cross-encoder model evaluates the relevance of all retrieved chunks and drops low-confidence data.", ACCENT_AMBER),
    ("LLM Synthesis", "The local 8B LLM receives the prompt + top documents and streams the answer back token-by-token.", ACCENT_RED)
]

for i, (title, desc, color) in enumerate(flow):
    top = Inches(3.0 + i * 0.9)
    badge = add_shape(slide, Inches(0.8), top, Inches(2.5), Inches(0.6), color)
    tf = add_text_box(slide, Inches(0.9), top + Inches(0.1), Inches(2.3), Inches(0.4))
    first_para(tf, title, font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    
    tf = add_text_box(slide, Inches(3.5), top, Inches(8.5), Inches(0.6))
    first_para(tf, desc, font_size=14, color=DARK_GRAY)

add_slide_number(slide, slide_num, TOTAL_SLIDES)

# Slide 3: Backend & API Layer
slide = next_slide()
add_bg(slide, WHITE)
tf = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8))
first_para(tf, "Backend & API Architecture", font_size=36, bold=True, color=BLACK)
add_divider(slide, Inches(0.8), Inches(1.2), Inches(2), ACCENT_TEAL)

tf = add_text_box(slide, Inches(0.8), Inches(2.0), Inches(11.5), Inches(0.5))
first_para(tf, "FastAPI provides the high-performance, asynchronous bridge connecting the AI to the UI.", font_size=16, color=DARK_GRAY)

card = add_shape(slide, Inches(0.8), Inches(2.8), Inches(11.5), Inches(3.5), LIGHT_GRAY, ACCENT_TEAL)

features_be = [
    ("Singleton Resource Management", "Database connections (Qdrant, Neo4j, SQLite) and heavy AI models are loaded once during application startup (lifespan events) to prevent out-of-memory crashes."),
    ("Streaming Responses", "Nginx proxy buffering is disabled. The FastAPI backend streams LLM outputs token-by-token, ensuring a snappy, low-latency chat experience even on slower GPUs."),
    ("Modular Routing", "API is divided logically into `/chat`, `/ingestion`, `/documents`, `/graph`, and `/compliance` endpoints for clean separation of concerns."),
]

for i, (title, desc) in enumerate(features_be):
    tf = add_text_box(slide, Inches(1.2), Inches(3.0 + i * 1.0), Inches(10.8), Inches(0.8))
    first_para(tf, f"• {title}:", font_size=16, bold=True, color=BLACK)
    set_para(tf, desc, font_size=14, color=DARK_GRAY)

add_slide_number(slide, slide_num, TOTAL_SLIDES)

# Slide 4: Knowledge Graph Construction
slide = next_slide()
add_bg(slide, WHITE)
tf = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8))
first_para(tf, "Knowledge Graph Construction", font_size=36, bold=True, color=BLACK)
add_divider(slide, Inches(0.8), Inches(1.2), Inches(2), ACCENT_BLUE)

tf = add_text_box(slide, Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.5))
first_para(tf, "Transforming flat documents into an interconnected web of assets.", font_size=20, bold=True, color=BLACK)
set_para(tf, "Cypher utilizes Qwen2.5-3B-Instruct during the ingestion phase for rapid Named Entity Recognition (NER). It identifies equipment tags, operating parameters, and regulatory norms, then dynamically creates nodes and edges in Neo4j.", font_size=16, color=DARK_GRAY, space_before=Pt(8))

card = add_shape(slide, Inches(0.8), Inches(4.0), Inches(11.5), Inches(2.0), LIGHT_GRAY, ACCENT_BLUE)
tf = add_text_box(slide, Inches(1.2), Inches(4.3), Inches(10.8), Inches(1.5))
first_para(tf, "Example Schema:", font_size=16, bold=True, color=ACCENT_BLUE)
set_para(tf, "(Document: Maintenance Report) -[MENTIONS]-> (Asset: Turbine 04)", font_size=16, bold=True, color=BLACK, space_before=Pt(8))
set_para(tf, "(Asset: Turbine 04) -[HAS_PARAMETER]-> (Parameter: 3000 RPM)", font_size=16, bold=True, color=BLACK, space_before=Pt(4))
set_para(tf, "This topology allows the LLM to deduce root causes across historically isolated incidents.", font_size=14, color=DARK_GRAY, space_before=Pt(8))

add_slide_number(slide, slide_num, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════
# SECTION 5 — BUSINESS ROADMAP
# ═══════════════════════════════════════════════════════════════════════
slide = next_slide()
add_section_header_slide(slide, "Business Roadmap", "The rollout strategy and value proposition for Cypher.", 5, ACCENT_TEAL)
add_slide_number(slide, slide_num, TOTAL_SLIDES)

# Slide 1: Phase 1 & 2
slide = next_slide()
add_bg(slide, WHITE)
tf = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8))
first_para(tf, "Implementation Roadmap — Phases 1 & 2", font_size=36, bold=True, color=BLACK)
add_divider(slide, Inches(0.8), Inches(1.2), Inches(2), ACCENT_TEAL)

card1 = add_shape(slide, Inches(0.8), Inches(2.0), Inches(11.5), Inches(2.0), LIGHT_GRAY, ACCENT_TEAL)
tf = add_text_box(slide, Inches(1.2), Inches(2.2), Inches(10.8), Inches(1.5))
first_para(tf, "Phase 1: Pilot & Core Integration (Months 1-3)", font_size=18, bold=True, color=ACCENT_TEAL)
set_para(tf, "• Deploy local edge-server in a single pilot plant.\n• Batch ingest 5 years of historical maintenance logs, OEM manuals, and P&IDs.\n• Roll out the 'Expert Knowledge Copilot' to the primary engineering team.", font_size=14, color=DARK_GRAY)

card2 = add_shape(slide, Inches(0.8), Inches(4.5), Inches(11.5), Inches(2.0), LIGHT_GRAY, ACCENT_BLUE)
tf = add_text_box(slide, Inches(1.2), Inches(4.7), Inches(10.8), Inches(1.5))
first_para(tf, "Phase 2: Advanced Agentic Capabilities (Months 4-6)", font_size=18, bold=True, color=ACCENT_BLUE)
set_para(tf, "• Activate the Maintenance Intelligence Engine for predictive RCA support.\n• Introduce Regulatory Compliance intelligence to auto-generate audit evidence.\n• Connect live telemetry to the knowledge graph.", font_size=14, color=DARK_GRAY)

add_slide_number(slide, slide_num, TOTAL_SLIDES)

# Slide 2: Phase 3 & ROI
slide = next_slide()
add_bg(slide, WHITE)
tf = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8))
first_para(tf, "Roadmap Phase 3 & ROI", font_size=36, bold=True, color=BLACK)
add_divider(slide, Inches(0.8), Inches(1.2), Inches(2), ACCENT_AMBER)

card1 = add_shape(slide, Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.8), LIGHT_GRAY, ACCENT_AMBER)
tf = add_text_box(slide, Inches(1.2), Inches(2.2), Inches(10.8), Inches(1.5))
first_para(tf, "Phase 3: Scale & Edge Mobility (Months 7-12)", font_size=18, bold=True, color=ACCENT_AMBER)
set_para(tf, "• Scale horizontally to remaining plants via synchronized distributed servers.\n• Deploy mobile/tablet interfaces for field technicians to query data on-site.\n• Potential AR (Augmented Reality) integrations for schematic overlays.", font_size=14, color=DARK_GRAY)

tf = add_text_box(slide, Inches(0.8), Inches(4.2), Inches(11.5), Inches(0.5))
first_para(tf, "Value Proposition & Return on Investment", font_size=20, bold=True, color=BLACK)

metrics = [
    ("35% → 5%", "Reduction in time spent searching for information.", ACCENT_TEAL),
    ("-15%", "Estimated decrease in unplanned downtime.", ACCENT_BLUE),
    ("100%", "Retention of institutional knowledge from retiring experts.", ACCENT_AMBER)
]

for i, (num, desc, color) in enumerate(metrics):
    left = Inches(0.8 + i * 4.0)
    card = add_shape(slide, left, Inches(5.0), Inches(3.6), Inches(1.5), WHITE, color)
    tf = add_text_box(slide, left + Inches(0.1), Inches(5.2), Inches(3.4), Inches(0.4))
    first_para(tf, num, font_size=28, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    tf = add_text_box(slide, left + Inches(0.1), Inches(5.8), Inches(3.4), Inches(0.6))
    first_para(tf, desc, font_size=12, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, slide_num, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════════════════
# SECTION 6 — RESULTS (Placeholders)
# ═══════════════════════════════════════════════════════════════════════
slide = next_slide()
add_section_header_slide(slide, "Results", "Showcasing the Cypher UI and capabilities.", 6, ACCENT_RED)
add_slide_number(slide, slide_num, TOTAL_SLIDES)

# Slide 1: Universal Search
slide = next_slide()
add_bg(slide, WHITE)
tf = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.8))
first_para(tf, "Universal Document Ingestion & Search", font_size=36, bold=True, color=BLACK)
add_divider(slide, Inches(0.8), Inches(1.2), Inches(2), ACCENT_BLUE)

ph_shape = add_shape(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(5.0), LIGHT_GRAY, ACCENT_BLUE)
tf = add_text_box(slide, Inches(0.8), Inches(4.0), Inches(11.5), Inches(1.0))
first_para(tf, "[ PLACEHOLDER: Screenshot of Document Ingestion/Search UI ]", font_size=20, bold=True, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, slide_num, TOTAL_SLIDES)

# Slide 2: Knowledge Graph Visualization
slide = next_slide()
add_bg(slide, WHITE)
tf = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.8))
first_para(tf, "Knowledge Graph Integration", font_size=36, bold=True, color=BLACK)
add_divider(slide, Inches(0.8), Inches(1.2), Inches(2), ACCENT_TEAL)

ph_shape = add_shape(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(5.0), LIGHT_GRAY, ACCENT_TEAL)
tf = add_text_box(slide, Inches(0.8), Inches(4.0), Inches(11.5), Inches(1.0))
first_para(tf, "[ PLACEHOLDER: Screenshot of Neo4j Graph Network/Visualization ]", font_size=20, bold=True, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, slide_num, TOTAL_SLIDES)

# Slide 3: Copilot
slide = next_slide()
add_bg(slide, WHITE)
tf = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.8))
first_para(tf, "Expert Knowledge Copilot", font_size=36, bold=True, color=BLACK)
add_divider(slide, Inches(0.8), Inches(1.2), Inches(2), ACCENT_AMBER)

ph_shape = add_shape(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(5.0), LIGHT_GRAY, ACCENT_AMBER)
tf = add_text_box(slide, Inches(0.8), Inches(4.0), Inches(11.5), Inches(1.0))
first_para(tf, "[ PLACEHOLDER: Screenshot of Chat UI answering a complex query with citations ]", font_size=20, bold=True, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, slide_num, TOTAL_SLIDES)

# Slide 4: Metrics
slide = next_slide()
add_bg(slide, WHITE)
tf = add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.8))
first_para(tf, "Performance & GPU Metrics", font_size=36, bold=True, color=BLACK)
add_divider(slide, Inches(0.8), Inches(1.2), Inches(2), ACCENT_RED)

ph_shape = add_shape(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(5.0), LIGHT_GRAY, ACCENT_RED)
tf = add_text_box(slide, Inches(0.8), Inches(4.0), Inches(11.5), Inches(1.0))
first_para(tf, "[ PLACEHOLDER: Screenshot of terminal logs or dashboard showing inference speed/VRAM ]", font_size=20, bold=True, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, slide_num, TOTAL_SLIDES)

# ═══════════════════════════════════════════════════════════════════════
# Save the Presentation
# ═══════════════════════════════════════════════════════════════════════
prs.save("/home/bw/CODES/Cypher/Cypher_Project_Presentation.pptx")
print("Presentation generated at /home/bw/CODES/Cypher/Cypher_Project_Presentation.pptx")
