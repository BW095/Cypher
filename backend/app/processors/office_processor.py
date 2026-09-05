"""
Office Document Processor — cloud deployment (no Docling).

Uses python-docx, python-pptx, and built-in readers.
"""

import os
import re
import traceback
from app.ingestion.canonical_document import CanonicalDocument


class OfficeProcessor:
    def process(self, file_path: str) -> CanonicalDocument:
        print(f"Processing Office Document: {file_path}")
        ext = os.path.splitext(file_path)[1].lower()

        handlers = {
            ".txt":  self._process_plaintext,
            ".html": self._process_html,
            ".htm":  self._process_html,
            ".docx": self._process_docx,
            ".doc":  self._process_docx,
            ".pptx": self._process_pptx,
            ".ppt":  self._process_pptx,
            ".odt":  self._process_plaintext,  # best-effort
        }
        handler = handlers.get(ext, self._process_plaintext)
        return handler(file_path)

    # ------------------------------------------------------------------
    def _process_plaintext(self, file_path: str) -> CanonicalDocument:
        try:
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                text = f.read()
            return CanonicalDocument(
                file_path=file_path, file_type="office_document",
                text=text, metadata={"processor": "direct_read"},
            )
        except Exception as e:
            return _error_doc(file_path, str(e), "direct_read_failed")

    def _process_html(self, file_path: str) -> CanonicalDocument:
        try:
            from bs4 import BeautifulSoup
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                soup = BeautifulSoup(f.read(), "lxml")
            text = soup.get_text(separator="\n")
            return CanonicalDocument(
                file_path=file_path, file_type="office_document",
                text=text, metadata={"processor": "beautifulsoup"},
            )
        except ImportError:
            # Fallback — strip tags with regex
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                raw = f.read()
            text = re.sub(r"<[^>]+>", " ", raw)
            text = re.sub(r"\s+", " ", text).strip()
            return CanonicalDocument(
                file_path=file_path, file_type="office_document",
                text=text, metadata={"processor": "html_regex"},
            )
        except Exception as e:
            return _error_doc(file_path, str(e), "html_failed")

    def _process_docx(self, file_path: str) -> CanonicalDocument:
        try:
            from docx import Document
            doc = Document(file_path)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            # Also extract tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip():
                        paragraphs.append(row_text)
            return CanonicalDocument(
                file_path=file_path, file_type="office_document",
                text="\n".join(paragraphs),
                metadata={"processor": "python-docx"},
            )
        except Exception as e:
            print(f"  python-docx failed: {e}")
            return _error_doc(file_path, str(e), "docx_failed")

    def _process_pptx(self, file_path: str) -> CanonicalDocument:
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            slides_text = []
            for i, slide in enumerate(prs.slides, 1):
                slide_parts = [f"[Slide {i}]"]
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_parts.append(shape.text.strip())
                slides_text.append("\n".join(slide_parts))
            return CanonicalDocument(
                file_path=file_path, file_type="office_document",
                text="\n\n".join(slides_text),
                metadata={"processor": "python-pptx"},
            )
        except Exception as e:
            print(f"  python-pptx failed: {e}")
            return _error_doc(file_path, str(e), "pptx_failed")


def _error_doc(file_path: str, err: str, processor: str) -> CanonicalDocument:
    return CanonicalDocument(
        file_path=file_path, file_type="office_document",
        text=f"[Processing failed: {err}]",
        metadata={"processor": processor},
    )